# %%
import os
from io import BufferedReader
from typing import Optional

# from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx
from loguru import logger
import magic
from transformers import AutoTokenizer, AutoModel
import torch
from typing import Dict, List, Optional, Tuple
import torch.nn.functional as F
from sentence_transformers import SentenceTransformer, util
import tiktoken
import sys
from utils.constant import  model_dir


csv.field_size_limit(sys.maxsize)

llm_tokenizer = tiktoken.get_encoding("cl100k_base")


def complete_sentence(text):
    last_dot_index = text.rfind(".")
    if last_dot_index != -1:
        cropped_text = text[: last_dot_index + 1]
        return cropped_text
    else:
        return text


# from models.models import Document, DocumentMetadata


# async def get_document_from_file(
#     file: UploadFile, metadata: DocumentMetadata
# ) -> Document:
#     extracted_text = await extract_text_from_form_file(file)

#     doc = Document(text=extracted_text, metadata=metadata)

#     return doc


# %%
def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    try:
        with open(filepath, "rb") as file:
            extracted_text = extract_text_from_file(file, mimetype)
    except Exception as e:
        logger.error(e)
        raise e

    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif (mimetype == "text/csv") or  (mimetype == "application/csv"):
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        raise ValueError("Unsupported file type: {}".format(mimetype))

    return extracted_text


# Extract text from a file based on its mimetype
def extract_text_from_form_all_file(
    folder_location, file_path_arr
):  # (file: UploadFile):
    """Return the text content of all files as array."""
    all_text_arr = []
    mime = magic.Magic(mime=True)

    for file_path in file_path_arr:
        file_loc = os.path.join(folder_location, file_path)
        mimetype = mime.from_file(file_loc)

        extracted_text = extract_text_from_filepath(
            file_loc, mimetype
        )  # (temp_file_path, mimetype)
        all_text_arr.extend([extracted_text])
    return all_text_arr


CHUNK_SIZE = 200  # The target size of each text chunk in tokens
MIN_CHUNK_SIZE_CHARS = 350  # The minimum size of each text chunk in characters
MIN_CHUNK_LENGTH_TO_EMBED = 10  # Discard chunks shorter than this
MAX_NUM_CHUNKS = 10000  # The maximum number of chunks to generate from a text


# %%
def get_files_chunks(tokenizer, text_arr, chunk_token_size):
    all_text_junks = []
    for text in text_arr:
        all_text_junks.extend(get_text_chunks(tokenizer, text, chunk_token_size))
    return all_text_junks


def get_text_chunks(
    tokenizer, text: str, chunk_token_size: Optional[int] = None
) -> List[str]:
    """
    Split a text into chunks of ~CHUNK_SIZE tokens, based on punctuation and newline boundaries.

    Args:
        text: The text to split into chunks.
        chunk_token_size: The target size of each chunk in tokens, or None to use the default CHUNK_SIZE.

    Returns:
        A list of text chunks, each of which is a string of ~CHUNK_SIZE tokens.
    """
    # Return an empty list if the text is empty or whitespace
    if not text or text.isspace():
        return []

    # Tokenize the text
    tokens = tokenizer.encode(
        text, disallowed_special=()
    )  # tokenizer.encode(text, add_special_tokens= False) # tokens = tokenizer.encode(text, disallowed_special=())
    # print(len(tokens))
    # return tokens
    # Initialize an empty list of chunks
    chunks = []

    # Use the provided chunk token size or the default one
    chunk_size = chunk_token_size or CHUNK_SIZE
    print("print chunk size:", chunk_size)

    # Initialize a counter for the number of chunks
    num_chunks = 0
    # Loop until all tokens are consumed
    while tokens and num_chunks < MAX_NUM_CHUNKS:
        # Take the first chunk_size tokens as a chunk
        chunk = tokens[:chunk_size]

        # Decode the chunk into text
        chunk_text = tokenizer.decode(
            chunk
        )  # tokenizer.decode(chunk, skip_special_tokens= True)

        # Skip the chunk if it is empty or whitespace
        if not chunk_text or chunk_text.isspace():
            # Remove the tokens corresponding to the chunk text from the remaining tokens
            tokens = tokens[len(chunk) :]
            # Continue to the next iteration of the loop
            continue

        # Find the last period or punctuation mark in the chunk
        last_punctuation = max(
            chunk_text.rfind("."),
            chunk_text.rfind("?"),
            chunk_text.rfind("!"),
            chunk_text.rfind("\n"),
        )

        # If there is a punctuation mark, and the last punctuation index is before MIN_CHUNK_SIZE_CHARS
        if last_punctuation != -1 and last_punctuation > MIN_CHUNK_SIZE_CHARS:
            # Truncate the chunk text at the punctuation mark
            chunk_text = chunk_text[: last_punctuation + 1]

        # Remove any newline characters and strip any leading or trailing whitespace
        chunk_text_to_append = chunk_text.replace("\n", " ").strip()

        if len(chunk_text_to_append) > MIN_CHUNK_LENGTH_TO_EMBED:
            # Append the chunk text to the list of chunks
            chunks.append(chunk_text_to_append)

        # Remove the tokens corresponding to the chunk text from the remaining tokens
        tokens = tokens[
            len(tokenizer.encode(chunk_text, disallowed_special=())) :
        ]  # len(tokenizer.encode(chunk_text, add_special_tokens= False)) :]

        # Increment the number of chunks
        num_chunks += 1

    # Handle the remaining tokens
    if tokens:
        remaining_text = (
            tokenizer.decode(chunk).replace("\n", " ").strip()
        )  # tokenizer.decode(tokens, skip_special_tokens= True).replace("\n", " ").strip()
        if len(remaining_text) > MIN_CHUNK_LENGTH_TO_EMBED:
            chunks.append(remaining_text)
    return chunks


# %%
class DocumentSematicSearch:
    def __init__(self, model_name, dir_name= model_dir):
        model_name = model_name #'hkunlp/instructor-base'#'sentence-transformers/all-mpnet-base-v2'
        self.tokenizer = AutoTokenizer.from_pretrained(model_name)
        self.model = AutoModel.from_pretrained(model_name)#.to("cuda")
        self.llm_tokenizer = llm_tokenizer
        self.dir_name = dir_name

    def check_model_exist(self, model_path):
        fpath = os.path.basename(os.path.normpath(model_path))
        if os.path.exists(f"{self.dir_name}{fpath}"):
            return f"{self.dir_name}{fpath}"
        # else:
        #     return model_path

         

    def mean_pooling(self, model_output, attention_mask):
        token_embeddings = model_output[
            0
        ]  # First element of model_output contains all token embeddings
        input_mask_expanded = (
            attention_mask.unsqueeze(-1).expand(token_embeddings.size()).float()
        )
        sum_embeddings = torch.sum(token_embeddings * input_mask_expanded, 1)
        sum_mask = torch.clamp(input_mask_expanded.sum(1), min=1e-9)
        return sum_embeddings / sum_mask

    def get_document_embedding(self, text):
        encoded_input = self.tokenizer(
            text, padding=True, truncation=True, return_tensors="pt"
        )
        encoded_input = encoded_input#.to("cuda")
        with torch.no_grad():
            model_output = self.model(**encoded_input)
            sentence_embeddings = self.mean_pooling(
                model_output, encoded_input["attention_mask"]
            )
        return sentence_embeddings#.to("cpu")

    def get_topk_result(self, query_embeddings, corpus_embeddings, k=5):
        result = util.semantic_search(
            query_embeddings, corpus_embeddings, score_function=util.cos_sim
        )
        return [[scores["corpus_id"], scores["score"]] for scores in result[0][:k]]

    # def get_document_index():


# %%
if __name__ == "__main__":
    text = extract_text_from_form_all_file(
        "/home/bioss/working_env/retriebal_solution/download_location/",
        ["Mohanlal - Wikipedia.pdf", "iCargo_Brochure.pdf"],
    )

    # tokenizer = AutoTokenizer.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")
    # model = AutoModel.from_pretrained("sentence-transformers/all-MiniLM-L6-v2")

    # document_token = get_text_chunks(tokenizer, text)
    doc_search_model = DocumentSematicSearch(
        "sentence-transformers/all-distilroberta-v1"
    )
    document_chunks = get_files_chunks(doc_search_model.tokenizer, text)
    doc_embd = doc_search_model.get_document_embedding(document_chunks)
    query_emb = doc_search_model.get_document_embedding("who is the father of mohanlal")
    topk_docs = doc_search_model.get_topk_result(query_emb, doc_embd)
    result = " ".join([document_chunks[doc_info[0]] for doc_info in topk_docs])


# %%

# %%


# %%

# def mean_pooling(model_output, attention_mask):
#     token_embeddings = model_output[0] #First element of model_output contains all token embeddings
#     input_mask_expanded = attention_mask.unsqueeze(-1).expand(token_embeddings.size()).float()
#     sum_embeddings = torch.sum(token_embeddings * input_mask_expanded, 1)
#     sum_mask = torch.clamp(input_mask_expanded.sum(1), min=1e-9)
#     return sum_embeddings / sum_mask
#     encoded_input = tokenizer(text, return_tensors='pt')
#     tokenizer.decode(encoded_input["input_ids"][0])
#     with torch.no_grad():
#         model_output = model(**encoded_input)
#         sentence_embeddings = mean_pooling(model_output, encoded_input['attention_mask'])


# %%
